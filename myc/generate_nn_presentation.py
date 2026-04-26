#!/usr/bin/env python3
from __future__ import annotations

import json
import sys
from pathlib import Path

VENDOR_DIR = Path(__file__).resolve().parent / ".vendor"
if str(VENDOR_DIR) not in sys.path:
    sys.path.insert(0, str(VENDOR_DIR))

import cairosvg
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from generate_nn_report import COMMON_RUN_DEFAULTS, PLOT_ITER_LIMIT, SUMMARY_JSON_PATH, fmt, load_algorithm_configs


OUT_DIR = Path(__file__).resolve().parent
ASSETS_DIR = OUT_DIR / "assets"
PNG_DIR = OUT_DIR / "ppt_assets"
PPT_PATH = OUT_DIR / "nn_experiment_report.pptx"
SCRIPT_PATH = OUT_DIR / "nn_experiment_report_script.md"
SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(246, 243, 237)
PAPER = RGBColor(255, 253, 249)
INK = RGBColor(29, 36, 51)
MUTED = RGBColor(93, 100, 114)
ACCENT = RGBColor(53, 80, 112)
ACCENT_2 = RGBColor(40, 114, 113)
GOOD = RGBColor(47, 133, 90)
WARN = RGBColor(192, 86, 33)
BAD = RGBColor(197, 48, 48)
CARD = RGBColor(243, 239, 232)

TITLE_FONT = "Microsoft YaHei"
BODY_FONT = "Microsoft YaHei"
MONO_FONT = "Consolas"


def main() -> None:
    data = json.loads(SUMMARY_JSON_PATH.read_text())
    runs = data["runs"]
    charts = data["charts"]
    algo_configs = load_algorithm_configs()

    by_short = {run["short_name"]: run for run in runs}
    solved = [run for run in runs if run["plot_metrics"]["first_one_peg_epoch"] is not None]
    best_final = max(run["plot_metrics"]["final_eval_reward"] for run in runs)
    best_final_runs = [run["short_name"] for run in runs if abs(run["plot_metrics"]["final_eval_reward"] - best_final) < 1e-9]
    fastest_epoch = min(solved, key=lambda item: item["plot_metrics"]["first_one_peg_epoch"])
    fastest_time = min(solved, key=lambda item: item["plot_metrics"]["time_to_one_peg_minutes"])
    biggest_gap = max(runs, key=lambda item: item["plot_metrics"]["eval_regression_gap"])
    lowest_entropy = sorted(runs, key=lambda item: item["plot_metrics"]["final_entropy"])[:3]
    highest_entropy = sorted(runs, key=lambda item: item["plot_metrics"]["final_entropy"], reverse=True)[:3]
    highest_critic = sorted(runs, key=lambda item: item["plot_metrics"]["final_critic_loss"], reverse=True)[:3]
    lowest_critic = sorted(runs, key=lambda item: item["plot_metrics"]["final_critic_loss"])[:3]

    PNG_DIR.mkdir(parents=True, exist_ok=True)
    pngs = {key: ensure_png(OUT_DIR / rel_path) for key, rel_path in charts.items()}

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    scripts: list[tuple[str, str]] = []
    total_slides = 13

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "FC / CNN / Transformer 训练日志分析汇报", "基于 checkpoints-and-logs/local，统一按前 3000 iter 口径")
    add_subtitle(
        slide,
        f"纳入 16 个带明确网络标签的实验，排除 4 个旧版 transformer_1d_old / transformer_2d_old 实验。",
        0.75,
        1.15,
        11.8,
        0.45,
        18,
    )
    add_stat_card(slide, 0.75, 1.85, 2.7, 1.35, "覆盖实验", "16", "8 个 A2C，8 个 PPO")
    add_stat_card(slide, 3.62, 1.85, 2.7, 1.35, "1-Peg 成功", "3", "A2C 1 个，PPO 2 个")
    add_stat_card(slide, 6.49, 1.85, 2.7, 1.35, "最佳最终评估", fmt(best_final, 3), "A2C · FC / PPO · CNN1 / PPO · CNN3")
    add_stat_card(
        slide,
        9.36,
        1.85,
        3.2,
        1.35,
        "最快到 1-Peg",
        fastest_epoch["short_name"],
        f"按 epoch {fastest_epoch['plot_metrics']['first_one_peg_epoch']}，按墙钟最快 {fastest_time['short_name']}",
    )
    add_quote_box(
        slide,
        0.75,
        3.55,
        11.8,
        2.2,
        "本次汇报重点",
        [
            f"所有图表和结论都只看前 {PLOT_ITER_LIMIT} iter。",
            "跨 FC / CNN / Transformer 的对比，本质上是“网络结构 + 该家族训练配置”的组合对比。",
            "更接近单变量消融的是 CNN1-5 内部对比，以及 Transformer 1D / 2D 的位置编码对比。",
        ],
    )
    scripts.append(
        (
            "Slide 1 封面",
            f"这份汇报分析的是本地日志目录下与 FC、CNN、Transformer 相关的 16 个实验。所有图表和主结论都统一限定在前 {PLOT_ITER_LIMIT} iter，同时完全排除了旧版的四个 transformer 旧实验。今天我会重点回答三个问题：第一，哪类方案在早期最有效；第二，哪些网络修改真的带来了收益；第三，critic loss 和策略熵这类训练指标，能不能帮助我们解释成功和失败。",
        )
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "一页结论", "先看最重要的判断，再回到各家族细节")
    add_bullet_panel(
        slide,
        0.75,
        1.1,
        12.0,
        5.6,
        [
            f"前 {PLOT_ITER_LIMIT} iter 内，最佳最终评估奖励是 {fmt(best_final, 3)}，由 {'、'.join(best_final_runs)} 并列拿到。",
            f"真正进入 1-peg 的只有 3 个实验，其中按评估 epoch 最快的是 {fastest_epoch['short_name']}，在 {fastest_epoch['plot_metrics']['first_one_peg_epoch']} epoch 达成；按墙钟时间最快的是 {fastest_time['short_name']}，只用了 {fmt(fastest_time['plot_metrics']['time_to_one_peg_minutes'], 1)} 分钟。",
            f"PPO 的优势主要集中在 CNN，尤其是 PPO · CNN1 和 PPO · CNN3；A2C 里最稳的是 FC。",
            f"Transformer 在保留实验里整体仍偏脆弱，尤其是 A2C · Transformer 1D，同时出现高梯度、低熵和较高 critic loss。",
            f"critic loss 和 entropy 都能解释稳定性，但都不能单独替代评估回报；最低 critic loss 并不一定代表学得最好。",
        ],
        font_size=20,
    )
    scripts.append(
        (
            "Slide 2 一页结论",
            f"如果先说结论，前 {PLOT_ITER_LIMIT} iter 的早期效率最强组合是 A2C · FC、PPO · CNN1 和 PPO · CNN3，它们窗口末评估奖励都达到 1.968。真正进入 1-peg 的只有 3 个实验，其中 PPO · CNN1 按样本效率最快，但 A2C · FC 按墙钟时间最省。再往下看，PPO 的优势主要体现在 CNN 上，A2C 最稳的基线是 FC。Transformer 这批实验还没有进入真正可用的阶段，尤其 A2C · Transformer 1D 同时暴露了梯度尖峰、熵塌缩和 critic loss 偏高的问题。",
        )
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "实验范围与口径", "先明确这份报告到底纳入了什么、又是怎么比较的")
    add_two_column_bullets(
        slide,
        "纳入范围",
        [
            "纳入所有目录名带 _fc、_cnn*、_transformer* 的实验。",
            "排除 1 个无架构后缀目录：A2C 2026-04-18 21-56-01。",
            "额外排除 4 个旧版 Transformer 实验：A2C/PPO + transformer_1d_old/transformer_2d_old。",
            "训练过程取自 training_history_full.csv；评估曲线取自 evaluation_metrics.csv。",
        ],
        "指标口径",
        [
            f"Best / Final evaluation reward 都限定在 iter <= {PLOT_ITER_LIMIT}。",
            f"1-peg epoch 指首次达到 pegs_left <= 1 且 epoch <= {PLOT_ITER_LIMIT}。",
            "Regression gap = best_eval - final_eval。",
            f"Final entropy / critic loss / gradient 都取窗口内最后一个训练点。",
        ],
    )
    scripts.append(
        (
            "Slide 3 实验范围与口径",
            f"这里先强调一下比较口径。第一，这次只纳入带明确网络标签的实验，也就是 FC、CNN 和 Transformer。第二，按要求排除了四个旧版 Transformer 实验，所以后面所有 Transformer 结论都只对应新批次的一维和二维位置编码。第三，所有指标都统一只看前 {PLOT_ITER_LIMIT} iter，这一点很重要，因为我们比较的是早期效率，而不是完整训练终局。",
        )
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "A2C / PPO 实验配置", "网络结构之外，算法训练器参数也决定了比较边界")
    add_two_column_bullets(
        slide,
        "公共运行参数",
        [
            f"n_envs = {COMMON_RUN_DEFAULTS['n_envs']}",
            f"n_steps = {COMMON_RUN_DEFAULTS['n_steps']}",
            f"device = {COMMON_RUN_DEFAULTS['device']}",
            f"enable_monitors = {str(COMMON_RUN_DEFAULTS['enable_monitors']).lower()}",
            "学习率默认来自各网络 YAML 的 optimizer.lr。",
        ],
        "A2C / PPO 训练器",
        [
            f"A2C: batch_size={algo_configs['A2C'].get('batch_size')}, n_iter={algo_configs['A2C'].get('n_iter')}, n_optim_steps={algo_configs['A2C'].get('n_optim_steps')}",
            f"A2C: n_games_train={algo_configs['A2C'].get('n_games_train')}, n_games_eval={algo_configs['A2C'].get('n_games_eval')}, seed={algo_configs['A2C'].get('seed')}",
            f"PPO: batch_size={algo_configs['PPO'].get('batch_size')}, n_iter={algo_configs['PPO'].get('n_iter')}, n_optim_steps={algo_configs['PPO'].get('n_optim_steps')}",
            f"PPO: clip_epsilon={algo_configs['PPO'].get('ppo', {}).get('clip_epsilon')}, entropy_coef={algo_configs['PPO'].get('ppo', {}).get('entropy_coef')}, value_loss_coef={algo_configs['PPO'].get('ppo', {}).get('value_loss_coef')}",
            f"PPO: max_grad_norm={algo_configs['PPO'].get('ppo', {}).get('max_grad_norm')}, gae_lambda={algo_configs['PPO'].get('ppo', {}).get('gae_lambda')}, discount={algo_configs['PPO'].get('ppo', {}).get('discount')}",
        ],
    )
    add_quote_box(
        slide,
        0.75,
        5.75,
        12.0,
        0.95,
        "解释约束",
        ["跨家族时不能把结果简单理解成“只改网络结构”；更准确的说法是“网络结构 + 该家族配置”的组合效果。"],
    )
    scripts.append(
        (
            "Slide 4 A2C / PPO 配置",
            "这里我把 A2C 和 PPO 的训练器参数单独列出来，是为了避免过度解读。虽然这批实验的主题是比较不同网络，但 FC、CNN 和 Transformer 的 YAML 并不只改了网络结构，同时也带着各自的学习率和损失配置。所以跨家族结论一定要保留这个前提；真正更接近单变量的，是 CNN 内部和 Transformer 1D/2D 内部的比较。",
        )
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "网络改动概览", "三类网络各自改动的轴并不一样")
    add_fitted_image(slide, pngs["architecture_fc"], 0.6, 1.25, 4.05, 3.0)
    add_fitted_image(slide, pngs["architecture_cnn"], 4.64, 1.25, 4.05, 3.0)
    add_fitted_image(slide, pngs["architecture_transformer"], 8.68, 1.25, 4.05, 3.0)
    add_bullet_panel(
        slide,
        0.75,
        4.55,
        12.0,
        2.0,
        [
            "FC: 7x7x3 展平后进入多层 GELU MLP，开销最小，但丢失 2D 局部结构先验。",
            "CNN1-5: 主要在卷积核大小、BatchNorm / GroupNorm、以及 PreAct + SE 残差增强之间做消融。",
            "Transformer: 核心差异是 1D 正弦位置编码和 2D 可学习行列位置编码。",
        ],
        font_size=18,
    )
    scripts.append(
        (
            "Slide 5 网络改动概览",
            "这一页的目的，是把三类网络到底改了什么讲清楚。FC 是最简单的全连接基线，优势是快，缺点是完全丢掉棋盘的空间结构。CNN 这一组实验，核心在测试卷积核大小、归一化方式以及更复杂的残差增强模块到底有没有用。Transformer 则主要是在比较一维和二维位置编码。后面所有结论，都应该对应到这些具体改动轴上。",
        )
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "全局表现", "先看早期窗口内谁跑得最好")
    add_fitted_image(slide, pngs["overall_final_eval"], 0.6, 1.15, 6.15, 4.8)
    add_fitted_image(slide, pngs["family_final_eval"], 6.9, 1.15, 5.85, 4.8)
    add_bullet_panel(
        slide,
        0.75,
        6.0,
        12.0,
        1.05,
        [
            f"窗口末最佳最终评估是 {fmt(best_final, 3)}，并列最优的是 {'、'.join(best_final_runs)}。",
            "PPO 的领先方案集中在 CNN；A2C 的最强单点来自 FC。",
        ],
        font_size=18,
    )
    scripts.append(
        (
            "Slide 6 全局表现",
            f"从全局图先看窗口末评估表现，前三个并列最优的分别是 A2C · FC、PPO · CNN1 和 PPO · CNN3，最终评估奖励都达到 {fmt(best_final, 3)}。这里可以明显看到一个分化：PPO 的优势集中在 CNN，A2C 的最强单点则来自 FC。也就是说，不同算法和不同网络家族之间，已经开始出现明显的匹配关系。",
        )
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "收敛效率与回退风险", "不能只看最高点，还要看是否能守住")
    add_fitted_image(slide, pngs["solved_efficiency"], 0.6, 1.2, 6.15, 4.7)
    add_fitted_image(slide, pngs["overall_regression_gap"], 6.9, 1.2, 5.85, 4.7)
    add_bullet_panel(
        slide,
        0.75,
        5.95,
        12.0,
        1.05,
        [
            f"最快到 1-peg 的是 {fastest_epoch['short_name']}，评估 epoch 为 {fastest_epoch['plot_metrics']['first_one_peg_epoch']}。",
            f"回退幅度最大的实验是 {biggest_gap['short_name']}，窗口内 regression gap 达到 {fmt(biggest_gap['plot_metrics']['eval_regression_gap'], 3)}。",
            f"保留实验里，Transformer 在 iter {PLOT_ITER_LIMIT} 前仍没有进入 1-peg。",
        ],
        font_size=18,
    )
    scripts.append(
        (
            "Slide 7 收敛效率与回退风险",
            f"这一页看两个稳定性指标。左边是到达 1-peg 的速度，右边是窗口内最好表现和窗口末表现之间的差，也就是 regression gap。能看到 PPO · CNN1 是最快进入 1-peg 的方案，而回退最严重的是 A2C · CNN3，它虽然中间一度接近 2 peg，但最后又掉回去了。这提醒我们，评估最高点不是唯一标准，能不能守住同样重要。",
        )
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "FC 家族", "A2C 最稳的基线来自 FC")
    add_fitted_image(slide, pngs["a2c_fc_eval"], 0.6, 1.2, 7.75, 5.15)
    add_bullet_panel(
        slide,
        8.55,
        1.35,
        4.2,
        4.9,
        [
            f"A2C · FC 在 {by_short['A2C · FC']['plot_metrics']['first_one_peg_epoch']} epoch 首次进入 1-peg，窗口末评估奖励 {fmt(by_short['A2C · FC']['plot_metrics']['final_eval_reward'], 3)}。",
            f"PPO · FC 停在 {fmt(by_short['PPO · FC']['plot_metrics']['final_eval_reward'], 3)}，稳定但没有冲到 1-peg。",
            f"FC 更像稳步爬升型基线，不是最激进的早期解法，但在 A2C 下表现非常健康。",
        ],
        font_size=18,
    )
    scripts.append(
        (
            "Slide 8 FC 家族",
            f"单独看 FC 家族，最值得注意的是 A2C · FC。它在 {by_short['A2C · FC']['plot_metrics']['first_one_peg_epoch']} epoch 首次进入 1-peg，并且最终把 1.968 的评估奖励稳稳守住了。相比之下，PPO · FC 更像是一个稳定但偏保守的方案，最后停在 0.968。我的理解是，FC 在 A2C 下形成了一个非常干净的基线：不是最快冲高的，但收敛形态很健康。",
        )
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "CNN 消融", "这一组是最接近单变量比较的网络实验")
    add_fitted_image(slide, pngs["a2c_cnn_eval"], 0.6, 1.15, 6.15, 3.7)
    add_fitted_image(slide, pngs["ppo_cnn_eval"], 6.9, 1.15, 5.85, 3.7)
    add_bullet_panel(
        slide,
        0.75,
        4.95,
        12.0,
        1.7,
        [
            "A2C-CNN 下，3x3 系列整体比 5x5 系列更稳，后者更容易在中期回退。",
            f"PPO-CNN 下，PPO · CNN1 和 PPO · CNN3 都在前 {PLOT_ITER_LIMIT} iter 内进入 1-peg，说明 5x5 核在 PPO 下并不吃亏。",
            f"PPO · CNN5 加入 PreAct + SE 后，评估从 {fmt(by_short['PPO · CNN5']['plot_metrics']['best_eval_reward'], 3)} 回退到 {fmt(by_short['PPO · CNN5']['plot_metrics']['final_eval_reward'], 3)}，当前证据不支持复杂残差增强能改善早期效率。",
        ],
        font_size=18,
    )
    scripts.append(
        (
            "Slide 9 CNN 消融",
            f"CNN 这一组是这批实验里最值得相信的消融，因为它们共享同一个卷积家族，只改了核大小、归一化和残差样式。结果上，A2C 更偏爱稳定的 3x3 方案，而 PPO 下最强的是 CNN1 和 CNN3，两者都在前 {PLOT_ITER_LIMIT} iter 进入 1-peg。反过来看，加入 PreAct 加 SE 的 CNN5 并没有带来收益，至少在早期窗口内没有。",
        )
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "Transformer 位置编码与稳定性", "新批次只保留 1D / 2D 的同窗比较")
    add_fitted_image(slide, pngs["a2c_transformer_eval"], 0.6, 1.1, 6.0, 3.2)
    add_fitted_image(slide, pngs["ppo_transformer_eval"], 6.85, 1.1, 5.9, 3.2)
    add_fitted_image(slide, pngs["transformer_gradients"], 7.15, 4.5, 5.55, 1.8)
    add_bullet_panel(
        slide,
        0.75,
        4.45,
        6.0,
        1.95,
        [
            f"A2C · Transformer 1D 的窗口末梯度达到 {fmt(by_short['A2C · Transformer 1D']['plot_metrics']['final_grad'], 3)}，熵只有 {fmt(by_short['A2C · Transformer 1D']['plot_metrics']['final_entropy'], 3)}。",
            f"A2C · Transformer 2D 的对应数值是 {fmt(by_short['A2C · Transformer 2D']['plot_metrics']['final_grad'], 3)} 和 {fmt(by_short['A2C · Transformer 2D']['plot_metrics']['final_entropy'], 3)}，稳定性明显更好。",
            f"但在 PPO 下，1D 和 2D 都没有在 iter {PLOT_ITER_LIMIT} 前进入 1-peg。",
        ],
        font_size=17,
    )
    scripts.append(
        (
            "Slide 10 Transformer 位置编码与稳定性",
            f"Transformer 这页的重点是稳定性，而不是单纯最终分数。A2C · Transformer 1D 的梯度在窗口末冲到 {fmt(by_short['A2C · Transformer 1D']['plot_metrics']['final_grad'], 3)}，同时熵掉到 {fmt(by_short['A2C · Transformer 1D']['plot_metrics']['final_entropy'], 3)}，这基本就是训练不稳的典型信号。换成 2D 位置编码后，梯度和熵都改善了，说明 2D 至少在稳定性上更合理。但即便如此，这一批 Transformer 在 PPO 下仍然没能进入 1-peg，所以它们目前更像是待继续调参的方向，而不是可直接替代 CNN 的成熟方案。",
        )
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "策略熵分析", "entropy 更接近“策略是否过早塌缩”")
    add_fitted_image(slide, pngs["entropy_final"], 0.6, 1.15, 8.2, 5.45)
    add_bullet_panel(
        slide,
        9.0,
        1.35,
        3.75,
        5.05,
        [
            f"最低熵的三个实验是 {lowest_entropy[0]['short_name']}、{lowest_entropy[1]['short_name']}、{lowest_entropy[2]['short_name']}。",
            f"前两个实验的窗口末评估奖励分别只有 {fmt(lowest_entropy[0]['plot_metrics']['final_eval_reward'], 3)} 和 {fmt(lowest_entropy[1]['plot_metrics']['final_eval_reward'], 3)}，典型塌缩。",
            f"熵最高的是 {highest_entropy[0]['short_name']}，达到 {fmt(highest_entropy[0]['plot_metrics']['final_entropy'], 3)}，但最终评估也只有 {fmt(highest_entropy[0]['plot_metrics']['final_eval_reward'], 3)}。",
            f"真正成功的 PPO · CNN1 / CNN3 处在中等熵区间，更像是“既保留探索、又没有完全塌缩”的平衡点。",
        ],
        font_size=17,
    )
    scripts.append(
        (
            "Slide 11 策略熵分析",
            f"这页专门看策略熵。最低熵的几个实验几乎都对应塌缩，尤其是 A2C · CNN1 和 A2C · CNN3，最后评估已经掉到 {fmt(lowest_entropy[0]['plot_metrics']['final_eval_reward'], 3)} 和 {fmt(lowest_entropy[1]['plot_metrics']['final_eval_reward'], 3)}。但反过来，熵也不是越高越好，最高熵的 PPO · CNN5 最终表现也不强。所以我更倾向于把熵理解成一个平衡指标：太低通常意味着策略已经过早确定，太高又说明策略还不够收敛，而真正成功的 PPO · CNN1 和 CNN3 恰好落在一个中等区间。",
        )
    )

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "Critic Loss 分析", "critic loss 能解释价值函数拟合状态，但不能单独代表效果")
    add_fitted_image(slide, pngs["critic_loss_final"], 0.6, 1.15, 8.2, 5.45)
    add_bullet_panel(
        slide,
        9.0,
        1.35,
        3.75,
        5.05,
        [
            f"最高 critic loss 来自 {highest_critic[0]['short_name']}，达到 {fmt(highest_critic[0]['plot_metrics']['final_critic_loss'], 3)}，这和它的梯度尖峰是一致的。",
            f"但最低 critic loss 的实验里也有 {lowest_critic[0]['short_name']} 和 {lowest_critic[1]['short_name']}，它们的最终评估分别只有 {fmt(lowest_critic[0]['plot_metrics']['final_eval_reward'], 3)} 和 {fmt(lowest_critic[1]['plot_metrics']['final_eval_reward'], 3)}。",
            f"{by_short['A2C · FC']['short_name']} 的 critic loss / entropy 为 {fmt(by_short['A2C · FC']['plot_metrics']['final_critic_loss'], 3)} / {fmt(by_short['A2C · FC']['plot_metrics']['final_entropy'], 3)}，属于更健康的组合。",
            "结论是：critic loss 更适合和评估回报、熵、梯度一起看，而不是单独排序。",
        ],
        font_size=17,
    )
    scripts.append(
        (
            "Slide 12 Critic Loss 分析",
            f"critic loss 这页最重要的结论是，它和最终效果并不是单调关系。最高 critic loss 的确对应 A2C · Transformer 1D 这种明显不稳定的实验，但最低 critic loss 的几个实验里，也包括最后表现很差的 A2C · CNN3 和 A2C · CNN1。这说明 critic loss 更像是价值函数当前拟合误差的一个侧面，它能帮助我们识别异常，但不能单独拿来决定哪个实验更好。更可靠的判断，还是要把 critic loss、entropy、gradient 和评估回报放在一起看。",
        )
    )

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    decorate_slide(slide, "效率与建议", "最后看训练成本，以及下一轮该怎么做")
    add_fitted_image(slide, pngs["runtime_per_epoch"], 0.6, 1.25, 7.75, 4.95)
    add_bullet_panel(
        slide,
        8.55,
        1.35,
        4.2,
        4.9,
        [
            f"A2C · FC 平均每个 epoch 只需 {fmt(by_short['A2C · FC']['plot_metrics']['seconds_per_epoch'], 3)} 秒，是最省时的强基线。",
            f"PPO · Transformer 1D 约 {fmt(by_short['PPO · Transformer 1D']['plot_metrics']['seconds_per_epoch'], 3)} 秒 / epoch，成本高且收益有限。",
            "如果目标是前 3000 iter 尽快拿到高评估表现，优先保留 PPO + CNN1 / CNN3。",
            "如果继续做 CNN 消融，建议围绕核大小和归一化再细化；PreAct + SE 目前没有证据支持。",
            "如果继续做 Transformer，先把稳定性当成第一优先级，再谈最终性能。",
        ],
        font_size=17,
    )
    scripts.append(
        (
            "Slide 13 效率与建议",
            f"最后看成本和行动建议。A2C · FC 是这批实验里性价比最高的强基线，单个 epoch 只要 {fmt(by_short['A2C · FC']['plot_metrics']['seconds_per_epoch'], 3)} 秒。相反，PPO · Transformer 1D 的单步开销已经到 {fmt(by_short['PPO · Transformer 1D']['plot_metrics']['seconds_per_epoch'], 3)} 秒，但窗口内效果并没有相应回报。所以下一轮如果追求前 {PLOT_ITER_LIMIT} iter 的早期性能，我会优先保留 PPO + CNN1 或 CNN3；如果要继续研究 Transformer，就先把训练稳定性问题解决掉。",
        )
    )

    for index, slide in enumerate(prs.slides, start=1):
        add_footer(slide, index, total_slides)

    prs.save(PPT_PATH)
    write_script(scripts)


def ensure_png(svg_path: Path) -> Path:
    png_path = PNG_DIR / f"{svg_path.stem}.png"
    cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), scale=2.0)
    return png_path


def decorate_slide(slide, title: str, kicker: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(SLIDE_W),
        Inches(0.88),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT
    banner.line.fill.background()

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        Inches(0.88),
        Inches(SLIDE_W),
        Inches(0.08),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_2
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(9.7), Inches(0.36))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = PAPER

    kicker_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.51), Inches(10.8), Inches(0.2))
    tf = kicker_box.text_frame
    p = tf.paragraphs[0]
    p.text = kicker
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(224, 231, 239)


def add_footer(slide, index: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.25), Inches(7.08), Inches(0.7), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = BODY_FONT
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_subtitle(slide, text: str, x: float, y: float, w: float, h: float, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = BODY_FONT
    p.font.size = Pt(font_size)
    p.font.color.rgb = INK


def add_stat_card(slide, x: float, y: float, w: float, h: float, label: str, metric: str, detail: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAPER
    shape.line.color.rgb = RGBColor(216, 209, 197)

    label_box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.1), Inches(w - 0.3), Inches(0.18))
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED

    metric_box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.32), Inches(w - 0.3), Inches(0.42))
    p = metric_box.text_frame.paragraphs[0]
    p.text = metric
    p.font.name = TITLE_FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    detail_box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.78), Inches(w - 0.3), Inches(0.36))
    tf = detail_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = detail
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = INK


def add_quote_box(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(216, 209, 197)

    head = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.36), Inches(0.22))
    p = head.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.42), Inches(w - 0.36), Inches(h - 0.52))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        if not p.text.startswith("• "):
            p.text = f"• {bullet}"
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = INK


def add_bullet_panel(slide, x: float, y: float, w: float, h: float, bullets: list[str], font_size: int = 18) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAPER
    shape.line.color.rgb = RGBColor(216, 209, 197)

    body = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.15), Inches(w - 0.36), Inches(h - 0.3))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK


def add_two_column_bullets(
    slide,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    add_panel_with_title(slide, 0.75, 1.25, 5.8, 5.65, left_title, left_bullets)
    add_panel_with_title(slide, 6.95, 1.25, 5.8, 5.65, right_title, right_bullets)


def add_panel_with_title(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAPER
    shape.line.color.rgb = RGBColor(216, 209, 197)

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.36), Inches(0.28))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.5), Inches(w - 0.36), Inches(h - 0.62))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(15)
        p.font.color.rgb = INK


def add_fitted_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        img_w, img_h = image.size
    box_w = Inches(w)
    box_h = Inches(h)
    ratio = min(box_w / img_w, box_h / img_h)
    draw_w = int(img_w * ratio)
    draw_h = int(img_h * ratio)
    left = Inches(x) + (box_w - draw_w) // 2
    top = Inches(y) + (box_h - draw_h) // 2

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = PAPER
    frame.line.color.rgb = RGBColor(216, 209, 197)

    slide.shapes.add_picture(str(path), left, top, width=draw_w, height=draw_h)


def write_script(scripts: list[tuple[str, str]]) -> None:
    parts = ["# NN 实验分析汇报逐字稿", "", f"这份逐字稿对应 [nn_experiment_report.pptx]({PPT_PATH.name})，建议控制在 10 到 15 分钟。", ""]
    for title, script in scripts:
        parts.append(f"## {title}")
        parts.append(script)
        parts.append("")
    SCRIPT_PATH.write_text("\n".join(parts))


if __name__ == "__main__":
    main()
